import streamlit as st
import pandas as pd
import copy
import io
import math

try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from docx import Document
    from docx.oxml import parse_xml, OxmlElement
    from docx.oxml.ns import nsdecls, qn
except ImportError as e:
    st.error(f"필수 라이브러리가 설치되지 않았습니다: {e}")
    st.stop()

st.set_page_config(page_title="대량 문서 자동 생성기", layout="centered")

st.title("📄 대량 양식 자동 생성기")
st.write("엑셀 데이터와 양식 템플릿(PPTX 또는 DOCX)을 업로드하세요.")
st.write("이름표, 태그 등 공통 양식의 대량 인쇄물 제작시 활용하세요.")

uploaded_excel = st.file_uploader("1. 엑셀 데이터 파일 (.xlsx)", type=["xlsx"])
uploaded_template = st.file_uploader("2. 템플릿 파일 (.pptx 또는 .docx)", type=["pptx", "docx"])


# ================= PPTX 관련 로직 =================
def get_all_ppt_shapes(shape_collection):
    shapes_list = []
    for shape in shape_collection:
        if shape.shape_type == 6 or hasattr(shape, 'shapes'):
            shapes_list.extend(get_all_ppt_shapes(shape.shapes))
        else:
            shapes_list.append(shape)
    return shapes_list


def copy_ppt_slide(prs, org_slide):
    copied_slide = prs.slides.add_slide(org_slide.slide_layout)
    for shape in org_slide.shapes:
        org_el = shape.element
        new_el = copy.deepcopy(org_el)
        copied_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    for value in org_slide.part.rels.values():
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.get_or_add(value.reltype, value._target)
    return copied_slide


def set_ppt_text_format(shape, text_val):
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    
    orig_size = p.font.size
    orig_name = p.font.name
    orig_bold = p.font.bold
    orig_color = p.font.color.rgb if p.font.color and hasattr(p.font.color, 'rgb') else None
    
    p.text = str(text_val)
    p.alignment = PP_ALIGN.CENTER
    
    if orig_size: p.font.size = orig_size
    if orig_name: p.font.name = orig_name
    if orig_bold: p.font.bold = orig_bold
    if orig_color: p.font.color.rgb = orig_color


def process_pptx(df, template_file):
    columns = [str(col).strip() for col in df.columns]
    prs = Presentation(template_file)
    first_slide = prs.slides[0]
    
    sample_shapes = get_all_ppt_shapes(first_slide.shapes)
    header_counts = {}
    for col in columns:
        count = sum(1 for s in sample_shapes if s.has_text_frame and s.text.strip() == col)
        if count > 0:
            header_counts[col] = count
            
    units_per_slide = max(header_counts.values()) if header_counts else 1
    total_slides = math.ceil(len(df) / units_per_slide)
    col_iterators = {col: iter(df[col]) for col in df.columns}
    
    for _ in range(total_slides):
        next_slide = copy_ppt_slide(prs, first_slide)
        all_shapes = get_all_ppt_shapes(next_slide.shapes)
        
        for shape in all_shapes:
            if shape.has_text_frame:
                clean_text = shape.text.strip()
                if clean_text in col_iterators:
                    try:
                        val = next(col_iterators[clean_text])
                        val = "" if pd.isna(val) else val
                        set_ppt_text_format(shape, val)
                    except StopIteration:
                        set_ppt_text_format(shape, "")
                        
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ================= DOCX 관련 로직 =================
def replace_text_in_runs(paragraph, key, val):
    """기존 폰트/크기/밑줄 서식을 유지하면서 정확히 키워드만 치환"""
    targets = [f"{{{{{key}}}}}", key]
    val_str = str(val)
    
    for target in targets:
        if target in paragraph.text:
            replaced = False
            for run in paragraph.runs:
                if target in run.text:
                    run.text = run.text.replace(target, val_str)
                    replaced = True
                    break
            
            if not replaced and paragraph.runs:
                # 런이 쪼개져 있는 경우 첫 런에 서식을 유지하며 치환
                paragraph.runs[0].text = paragraph.text.replace(target, val_str)
                for run in paragraph.runs[1:]:
                    run.text = ""


def process_docx(df, template_file):
    template_bytes = template_file.getvalue()
    columns = [str(col).strip() for col in df.columns]
    
    # 1. 템플릿 1장 안에 각 컬럼명이 몇 개씩 있는지 카운트
    sample_doc = Document(io.BytesIO(template_bytes))
    header_counts = {col: 0 for col in columns}
    
    for p in sample_doc.paragraphs:
        for col in columns:
            if col in p.text or f"{{{{{col}}}}}" in p.text:
                header_counts[col] += 1
                
    for table in sample_doc.tables:
        for r in table.rows:
            for cell in r.cells:
                for p in cell.paragraphs:
                    for col in columns:
                        if col in p.text or f"{{{{{col}}}}}" in p.text:
                            header_counts[col] += 1

    units_per_page = max(header_counts.values()) if any(header_counts.values()) else 1
    total_pages = math.ceil(len(df) / units_per_page)
    
    col_iterators = {col: iter(df[col]) for col in df.columns}
    
    def fill_doc(doc):
        """문서 내의 키워드를 순차 데이터로 치환"""
        for p in doc.paragraphs:
            for col in columns:
                if col in p.text or f"{{{{{col}}}}}" in p.text:
                    try:
                        val = next(col_iterators[col])
                        val = "" if pd.isna(val) else val
                    except StopIteration:
                        val = ""
                    replace_text_in_runs(p, col, val)
                    
        for table in doc.tables:
            for r in table.rows:
                for cell in r.cells:
                    for p in cell.paragraphs:
                        for col in columns:
                            if col in p.text or f"{{{{{col}}}}}" in p.text:
                                try:
                                    val = next(col_iterators[col])
                                    val = "" if pd.isna(val) else val
                                except StopIteration:
                                    val = ""
                                replace_text_in_runs(p, col, val)

    # 마스터 문서 초기화
    master_doc = Document(io.BytesIO(template_bytes))
    fill_doc(master_doc)
    
    # 2번째 페이지부터 결합할 때 페이지 넘김(pageBreakBefore) 속성을 첫 요소에 직접 적용
    for page_num in range(1, total_pages):
        page_doc = Document(io.BytesIO(template_bytes))
        fill_doc(page_doc)
        
        # page_doc의 본문 블록들 추출 (sectPr 제외)
        body_elements = [el for el in page_doc.element.body if not el.tag.endswith('sectPr')]
        
        if body_elements:
            first_el = body_elements[0]
            # 첫 요소가 단락(<w:p>)이면 해당 단락에 pageBreakBefore를 줘서 빈 줄 없이 깔끔하게 다음 페이지로 시작
            if first_el.tag.endswith('p'):
                pPr = first_el.find(qn('w:pPr'))
                if pPr is None:
                    pPr = OxmlElement('w:pPr')
                    first_el.insert(0, pPr)
                pPr.append(parse_xml(r'<w:pageBreakBefore %s/>' % nsdecls('w')))
            else:
                # 표(<w:tbl>)로 시작하는 경우 단락 구분자 추가
                br_p = parse_xml(r'<w:p %s><w:r><w:br w:type="page"/></w:r></w:p>' % nsdecls('w'))
                master_doc.element.body.append(br_p)
                
            for el in body_elements:
                master_doc.element.body.append(copy.deepcopy(el))
                
    output = io.BytesIO()
    master_doc.save(output)
    output.seek(0)
    return output


# ================= 메인 UI =================
if uploaded_excel and uploaded_template:
    try:
        df = pd.read_excel(uploaded_excel)
        columns = [str(col).strip() for col in df.columns]
        file_type = "pptx" if uploaded_template.name.endswith(".pptx") else "docx"
        
        st.success(f"엑셀 로드 완료: 총 {len(df)}개 행 / 감지된 헤더: {', '.join(columns)}")
        
        if st.button("🚀 대량 문서 생성 시작"):
            with st.spinner(f"{file_type.upper()} 문서를 생성하는 중..."):
                if file_type == "pptx":
                    output_stream = process_pptx(df, uploaded_template)
                    file_name = "생성결과.pptx"
                    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                else:
                    output_stream = process_docx(df, uploaded_template)
                    file_name = "생성결과.docx"
                    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    
                st.success("✨ 생성이 완료되었습니다!")
                st.download_button(
                    label=f"📥 완성된 {file_type.upper()} 파일 다운로드",
                    data=output_stream,
                    file_name=file_name,
                    mime=mime
                )
    except Exception as e:
        st.error(f"처리 중 오류 발생: {e}")
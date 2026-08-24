import streamlit as st
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import pandas as pd
import copy
import io
import math

st.set_page_config(page_title="대량 PPT 자동 생성기", layout="centered")

st.title("📄 대량 PPT 양식 자동 생성기")
st.write("엑셀 데이터와 PPT 템플릿을 업로드하면 엑셀의 헤더(열 이름)와 일치하는 PPT 텍스트를 자동으로 치환하여 완성합니다.")

# 파일 업로드 영역
uploaded_excel = st.file_uploader("1. 엑셀 파일 업로드 (.xlsx)", type=["xlsx"])
uploaded_pptx = st.file_uploader("2. PPT 템플릿 업로드 (.pptx)", type=["pptx"])


def get_all_shapes(shape_collection):
    """그룹 도형 내부까지 포함한 모든 단일 도형 리스트 반환"""
    shapes_list = []
    for shape in shape_collection:
        if shape.shape_type == 6 or hasattr(shape, 'shapes'):
            shapes_list.extend(get_all_shapes(shape.shapes))
        else:
            shapes_list.append(shape)
    return shapes_list


def copy_slide(prs, org_slide):
    """슬라이드 복제 함수"""
    copied_slide = prs.slides.add_slide(org_slide.slide_layout)
    for shape in org_slide.shapes:
        org_el = shape.element
        new_el = copy.deepcopy(org_el)
        copied_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    for value in org_slide.part.rels.values():
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.get_or_add(
                value.reltype,
                value._target
            )
    return copied_slide


def set_text_format(shape, text_val):
    """기존 글상자 서식을 유지하면서 텍스트 및 가운데 정렬 적용"""
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p = tf.paragraphs[0]
    
    # 기존 서식 보존
    original_size = p.font.size
    original_name = p.font.name
    original_bold = p.font.bold
    original_color = p.font.color.rgb if p.font.color and hasattr(p.font.color, 'rgb') else None
    
    # 값 입력 및 정렬
    p.text = str(text_val)
    p.alignment = PP_ALIGN.CENTER
    
    # 서식 복원
    if original_size is not None:
        p.font.size = original_size
    if original_name is not None:
        p.font.name = original_name
    if original_bold is not None:
        p.font.bold = original_bold
    if original_color is not None:
        p.font.color.rgb = original_color


if uploaded_excel and uploaded_pptx:
    # 데이터 로드
    df = pd.read_excel(uploaded_excel)
    columns = [str(col).strip() for col in df.columns]
    
    st.success(f"엑셀 로드 완료: 총 {len(df)}개의 데이터 행 감지됨")
    st.write(f"**감지된 열 헤더:** `{', '.join(columns)}`")
    
    if st.button("🚀 PPT 생성 시작"):
        with st.spinner("PPT를 생성하는 중입니다..."):
            prs = Presentation(uploaded_pptx)
            first_slide = prs.slides[0]
            
            # 첫 번째 슬라이드에서 엑셀 헤더와 매칭되는 텍스트 상자 개수 카운트 (슬라이드당 양식 개수 파악)
            sample_shapes = get_all_shapes(first_slide.shapes)
            header_counts = {}
            for col in columns:
                count = sum(1 for s in sample_shapes if s.has_text_frame and s.text.strip() == col)
                if count > 0:
                    header_counts[col] = count
            
            # 슬라이드 한 장당 몇 세트의 양식이 들어있는지 확인 (예: 2개씩 배치된 경우 units_per_slide = 2)
            units_per_slide = max(header_counts.values()) if header_counts else 1
            total_slides_needed = math.ceil(len(df) / units_per_slide)
            
            # 엑셀의 각 열 데이터를 iterator로 준비
            col_iterators = {col: iter(df[col]) for col in df.columns}
            
            # 필요한 슬라이드만큼 복제 및 데이터 치환
            for _ in range(total_slides_needed):
                next_slide = copy_slide(prs, first_slide)
                all_shapes = get_all_shapes(next_slide.shapes)
                
                for shape in all_shapes:
                    if shape.has_text_frame:
                        clean_text = shape.text.strip()
                        # 텍스트 상자의 내용이 엑셀 컬럼명과 일치하는 경우
                        if clean_text in col_iterators:
                            try:
                                val = next(col_iterators[clean_text])
                                # NaN 값 처리
                                if pd.isna(val):
                                    val = ""
                                set_text_format(shape, val)
                            except StopIteration:
                                # 데이터가 부족할 경우 빈칸 처리
                                set_text_format(shape, "")
            
            # 첫 번째 템플릿 원본 슬라이드 삭제 (필요한 경우)
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
            
            # 메모리 버퍼에 저장
            output_stream = io.BytesIO()
            prs.save(output_stream)
            output_stream.seek(0)
            
            st.success("✨ 생성이 완료되었습니다!")
            
            # 다운로드 버튼 제공
            st.download_button(
                label="📥 완성된 PPT 파일 다운로드",
                data=output_stream,
                file_name="생성결과.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )